from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import Dict
from fastapi.middleware.cors import CORSMiddleware
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from concurrent.futures import ThreadPoolExecutor
import time

from langchain_community.document_loaders import UnstructuredPDFLoader, TextLoader, DirectoryLoader
from langchain_community.document_loaders import UnstructuredWordDocumentLoader
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_huggingface.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.document_loaders import PyMuPDFLoader
from langchain.chains import MapReduceDocumentsChain, ReduceDocumentsChain
from langchain.chains.llm import LLMChain
from langchain_core.prompts import PromptTemplate
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
from langchain_text_splitters import CharacterTextSplitter

app = FastAPI()

executor = ThreadPoolExecutor(max_workers=10)

def background_task(text):

    loader = DirectoryLoader('./uploads/', show_progress=True)
    filename = "Summary.txt"
    # with open(filename, 'w', encoding='utf-8') as file:
    #     file.write(text)
    # loader = TextLoader(filename)
    data = loader.load()
    text_splitter = CharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=1500, chunk_overlap=70
    )
    # if os.path.isfile(filename):
    #     os.remove(filename)
    split_docs = text_splitter.split_documents(data)

    result = map_reduce_chain.invoke(split_docs)

    with open(filename, 'w', encoding='utf-8') as file:
        file.write(result["output_text"])
    
    print("Processed Summary")
    return

# Allow all origins, or specify allowed origins in a list
origins = [
    "http://localhost:3000",  # React frontend
    "http://localhost:8000",  # FastAPI backend, if accessed via browser
    "http://localhost:5173"
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class Item(BaseModel):
    data: Dict
class Question(BaseModel):
    question: str

# Create the uploads directory if it doesn't exist
if not os.path.exists("uploads"):
    os.makedirs("uploads")


model_name = "sentence-transformers/all-MiniLM-L6-v2"
model_kwargs = {'device': 'cpu'}
encode_kwargs = {'normalize_embeddings': False}
embeddings = HuggingFaceEmbeddings(
            model_name=model_name,
            model_kwargs=model_kwargs,
            encode_kwargs=encode_kwargs
            )

llm = ChatGoogleGenerativeAI(
        model="gemini-1.5-pro",
        temperature=0.5,
        max_tokens=None,
        timeout=None,
        max_retries=2,
        google_api_key="AIzaSyByAULS9YrPUqkrai_ZQn5-PPaOxBaTpcU"
        
    )

map_template = """The following is a set of documents
{docs}
Based on this list of docs, please identify the main themes 
Helpful Answer:"""
map_prompt = PromptTemplate.from_template(map_template)
map_chain = LLMChain(llm=llm, prompt=map_prompt)

reduce_template = """The following is set of summaries:
{docs}
Take these and distill it into a final, consolidated summary of the main themes. 
Helpful Answer:"""
reduce_prompt = PromptTemplate.from_template(reduce_template)

# Run chain
reduce_chain = LLMChain(llm=llm, prompt=reduce_prompt)

# Takes a list of documents, combines them into a single string, and passes this to an LLMChain
combine_documents_chain = StuffDocumentsChain(
    llm_chain=reduce_chain, document_variable_name="docs"
)

# Combines and iteratively reduces the mapped documents
reduce_documents_chain = ReduceDocumentsChain(
    # This is final chain that is called.
    combine_documents_chain=combine_documents_chain,
    # If documents exceed context for `StuffDocumentsChain`
    collapse_documents_chain=combine_documents_chain,
    # The maximum number of tokens to group documents into.
    token_max=4000,
)


# Combining documents by mapping a chain over them, then combining results
map_reduce_chain = MapReduceDocumentsChain(
    # Map chain
    llm_chain=map_chain,
    # Reduce chain
    reduce_documents_chain=reduce_documents_chain,
    # The variable name in the llm_chain to put the documents in
    document_variable_name="docs",
    # Return the results of the map steps in the output
    return_intermediate_steps=False,
)

text_splitter = CharacterTextSplitter(
    chunk_size=1000, chunk_overlap=0
)

def get_pdf_text(pdf_file):
    text = ""
    pdf_reader = PdfReader(pdf_file)
    for page_num, page in enumerate(pdf_reader.pages, 1):
        text += f"Page {page_num}:\n"
        text += page.extract_text()
        text += "\n\n"  # Add space between pages
    return text

def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    text = ""

    for slide_num, slide in enumerate(presentation.slides, 1):
        text += f"Slide {slide_num}:\n"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += f"{paragraph.text}\n"
        
        text += "\n"  # Add a separator between slides
    return text

def get_response(question):
    db = FAISS.load_local("vectorstore", embeddings, allow_dangerous_deserialization=True)
    system_prompt = (
        "Use the given context to answer the question. "
        "If you don't know the answer, say you don't know. "
        "Context: {context}"
    )
    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system_prompt),
            ("human", "{input}"),
        ]
    )

    retriever = db.as_retriever()

    question_answer_chain = create_stuff_documents_chain(llm, prompt)
    chain = create_retrieval_chain(retriever, question_answer_chain)

    response = chain.invoke({"input": question})
    return response['answer']

def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    text = ""
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            text += f"\n\n{para.text}\n\n"  # Add extra space around headings
        else:
            text += para.text + " "  # Combine lines into paragraphs

    # Ensure paragraphs are separated properly
    paragraphs = text.split("\n\n")
    formatted_text = "\n\n".join(para.strip() for para in paragraphs if para.strip())

    for table in doc.tables:
        formatted_text += "\n\n"
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    formatted_text += cell.text + "\t"
            formatted_text += "\n"
        formatted_text += "\n"

    return formatted_text

def delete_files_in_folder(folder_path):
    # Iterate over all the items in the folder
    for item in os.listdir(folder_path):
        item_path = os.path.join(folder_path, item)
        
        # Check if the item is a file
        if os.path.isfile(item_path):
            os.remove(item_path)  # Remove the file


@app.get("/")
async def read_root():
    return {"message": "Hello, World!"}

@app.get("/get-summary")
async def get_summary():

    filename = "Summary.txt"
    if os.path.isfile(filename):
        with open(filename, 'r') as file:
            content = file.read()
    else:
        content = "Please Request Summary after Sometime, it is currently being Processed. Thank you for your Patience"

    return JSONResponse(status_code=200, content={"summary": content})
        

@app.post("/upload-single/")
async def upload_file(file: UploadFile = File(...), data: str = Form(...)):
    try:
        # Save the file to the uploads directory
        file_location = f"uploads/{file.filename}"
        with open(file_location, "wb+") as file_object:
            file_object.write(file.file.read())
        print(data)
        # Process the JSON data if needed   
        # You can use json.loads(data) if you need to parse it

        return JSONResponse(status_code=200, content={"message": "File uploaded successfully"})
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": f"An error occurred: {str(e)}"})

@app.post("/upload-many/")  # Increase timeout to 1 hour
async def upload_files(files: list[UploadFile] = File(...), data: str = Form(...), background_tasks: BackgroundTasks = BackgroundTasks()):
    text_final = ""
    print(data)
    delete_files_in_folder("./uploads/")
    if os.path.isfile("Summary.txt"):
        os.remove("Summary.txt")
    try:
        for file in files:
            file_location = f"uploads/{file.filename}"
            with open(file_location, "wb+") as file_object:
                file_object.write(file.file.read())
            
            # Perform text extraction based on file type
            if file.filename.endswith('.pdf'):
                text = get_pdf_text(file_location).strip()
                text_final += text
            elif file.filename.endswith('.pptx'):
                text = extract_text_from_pptx(file_location).strip()
                text_final += text
            elif file.filename.endswith('.docx'):
                text = extract_text_from_docx(file_location).strip()
                text_final += text
            
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=900, chunk_overlap=100)
            texts = text_splitter.split_text(text_final)

            db = FAISS.from_texts(texts, embeddings)
            db.save_local("vectorstore")

        background_tasks.add_task(background_task, text_final)
            

        # Process the JSON data if needed
        # You can use json.loads(data) if you need to parse it

        return JSONResponse(status_code=200, content={"message": "Files uploaded and processed successfully", "extracted_text": text_final})
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": f"An error occurred: {str(e)}"})
    
@app.post("/ask-question/")
async def ask_question(question: Question):
    try:
        # Process the question here
        response_text = get_response(question.question)
        
        # Return the response as a JSON object
        return {"response": response_text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An error occurred: {str(e)}")
    
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="localhost", port=5000)
